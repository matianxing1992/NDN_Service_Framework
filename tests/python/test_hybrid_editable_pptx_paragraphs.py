#!/usr/bin/env python3

import importlib.util
from pathlib import Path
import sys
import unittest

from pptx import Presentation
from pptx.util import Inches


REPO_ROOT = Path(__file__).resolve().parents[2]
SLIDES = REPO_ROOT / "docs" / "PAPER" / "proposal-defense" / "slides"
sys.path.insert(0, str(SLIDES))
SPEC = importlib.util.spec_from_file_location(
    "hybrid_editable_pptx", SLIDES / "generate_hybrid_editable_pptx.py"
)
hybrid = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = hybrid
SPEC.loader.exec_module(hybrid)


FONT = hybrid.FontSpec(
    source_id="0",
    size=11.0,
    source_family="LMSans10-Regular",
    mapped_family="Arial",
    rgb=(0, 0, 0),
    family_bold=False,
    family_italic=False,
)
BULLET_FONT = hybrid.FontSpec(
    source_id="1",
    size=11.0,
    source_family="MSAM10",
    mapped_family="Arial",
    rgb=(26, 62, 122),
    family_bold=False,
    family_italic=False,
)
MONO_FONT = hybrid.FontSpec(
    source_id="7",
    size=11.0,
    source_family="LMMono8-Regular",
    mapped_family="Latin Modern Mono",
    rgb=(0, 0, 0),
    family_bold=False,
    family_italic=False,
)


def span(span_id, text, left, top, width, *, font=FONT):
    return hybrid.TextSpan(
        span_id=span_id,
        page=1,
        page_w=680.0,
        page_h=382.0,
        left=float(left),
        top=float(top),
        width=float(width),
        height=10.6,
        font=font,
        inline_runs=(hybrid.InlineTextRun(text, False, False),),
    )


def plan_for(spans, *, page_number=2):
    page = hybrid.PdfPage(page_number, 680.0, 382.0, tuple(spans))
    return page, hybrid.plan_conversion(hybrid.PdfDocument((page,)))


class HybridEditableParagraphTests(unittest.TestCase):
    def test_centered_multiline_label_keeps_hard_line_break(self):
        page = hybrid.PdfPage(
            44,
            680.0,
            382.0,
            (
                span(1, "Instructions for each Provider", 200, 20, 190),
                span(2, "task + model part + device", 215, 34.2, 160),
            ),
        )
        group = hybrid.TextGroup(
            1, 44, page.spans, line_start_span_ids=(2,), alignment="center"
        )
        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        hybrid.add_text_group(slide, page, group)
        self.assertEqual(
            slide.shapes[-1].text,
            "Instructions for each Provider\vtask + model part + device",
        )
        self.assertEqual(hybrid.expected_group_text(group), slide.shapes[-1].text)

    def test_semantic_space_is_prefixed_to_visible_run(self):
        page, plan = plan_for(
            [
                span(1, "Proposed Work Status and", 12, 20, 165),
                span(2, "Timeline", 181, 20, 52),
            ]
        )
        group = plan.groups_by_page[page.number][0]
        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        hybrid.add_text_group(slide, page, group)
        paragraph = slide.shapes[-1].text_frame.paragraphs[0]
        self.assertEqual(
            [run.text for run in paragraph.runs],
            ["Proposed Work Status and", " Timeline"],
        )
        self.assertFalse(any(run.text.isspace() for run in paragraph.runs))
        self.assertEqual(slide.shapes[-1].text, "Proposed Work Status and Timeline")

    def test_monospace_code_lines_remain_separate_nonwrapping_textboxes(self):
        page, plan = plan_for(
            [
                span(1, "provider /UAVNET/drone1", 52, 20, 146, font=MONO_FONT),
                span(2, "allow /FlightControl/Takeoff", 69, 34.2, 178, font=MONO_FONT),
                span(3, "allow /FlightControl/Land", 69, 48.4, 159, font=MONO_FONT),
            ]
        )
        groups = plan.groups_by_page[page.number]
        self.assertEqual(
            [hybrid.expected_group_text(group) for group in groups],
            [
                "provider /UAVNET/drone1",
                "allow /FlightControl/Takeoff",
                "allow /FlightControl/Land",
            ],
        )
        self.assertTrue(all(not group.line_start_span_ids for group in groups))

        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for group in groups:
            hybrid.add_text_group(slide, page, group)
        self.assertTrue(
            all(shape.text_frame.word_wrap is False for shape in slide.shapes)
        )

    def test_wrapped_visual_lines_become_one_paragraph_group(self):
        page, plan = plan_for(
            [
                span(1, "A paragraph that was split by the PDF", 100, 20, 230),
                span(2, "continues on the next visual line.", 100, 34.2, 190),
            ]
        )
        groups = plan.groups_by_page[page.number]
        self.assertEqual(len(groups), 1)
        self.assertEqual(groups[0].line_start_span_ids, (2,))
        self.assertEqual(
            hybrid.expected_group_text(groups[0]),
            "A paragraph that was split by the PDF continues on the next visual line.",
        )

    def test_each_bullet_starts_a_new_wrapping_textbox(self):
        page, plan = plan_for(
            [
                span(1, "I", 80, 20, 9, font=BULLET_FONT),
                span(2, "First bullet has a long first line", 100, 20, 210),
                span(3, "and one continuation line.", 100, 34.2, 160),
                span(4, "I", 80, 48.4, 9, font=BULLET_FONT),
                span(5, "Second bullet starts here.", 100, 48.4, 170),
            ]
        )
        groups = plan.groups_by_page[page.number]
        texts = [hybrid.expected_group_text(group) for group in groups]
        self.assertIn(
            "First bullet has a long first line and one continuation line.", texts
        )
        self.assertIn("Second bullet starts here.", texts)
        self.assertEqual(len(plan.bullets_by_page[page.number]), 2)

    def test_table_rows_do_not_merge_but_wrapped_cells_do(self):
        page, plan = plan_for(
            [
                span(1, "gRPC-1", 100, 20, 45),
                span(2, "one fixed target;", 200, 20, 100),
                span(3, "no failover", 200, 34.2, 70),
                span(4, "NSC-1", 100, 48.4, 40),
                span(5, "one fixed Provider", 200, 48.4, 110),
                span(6, "prefix; no failover", 200, 62.6, 110),
            ]
        )
        texts = [
            hybrid.expected_group_text(group)
            for group in plan.groups_by_page[page.number]
        ]
        self.assertIn("one fixed target; no failover", texts)
        self.assertIn("one fixed Provider prefix; no failover", texts)
        self.assertNotIn(
            "one fixed target; no failover one fixed Provider prefix; no failover",
            texts,
        )

    def test_rejected_table_row_cannot_attach_to_an_older_cell(self):
        page, plan = plan_for(
            [
                span(1, "gRPC-1", 100, 20, 45),
                span(2, "one", 200, 20, 18),
                span(3, "fixed", 228, 20, 24),
                span(4, "target;", 263, 20, 34),
                span(5, "no failover", 200, 34.2, 70),
                span(6, "NSC-1", 100, 48.4, 40),
                span(7, "one fixed Provider", 200, 48.4, 110),
                span(8, "prefix; no failover", 200, 62.6, 110),
                span(9, "NDNSF", 100, 76.8, 48),
                span(10, "four-provider", 200, 76.8, 80),
                span(11, "ACK", 200, 91.0, 25),
                span(12, "collection;", 246, 91.0, 54),
                span(13, "FirstResponding", 200, 105.2, 96),
            ]
        )
        texts = [
            hybrid.expected_group_text(group)
            for group in plan.groups_by_page[page.number]
        ]
        self.assertIn("one fixed target; no failover", texts)
        self.assertIn("one fixed Provider prefix; no failover", texts)
        self.assertIn("four-provider ACK collection; FirstResponding", texts)
        self.assertNotIn("one no failover", texts)
        self.assertNotIn("ACK FirstResponding", texts)

    def test_bold_label_and_same_line_body_stay_in_wrapped_paragraph(self):
        regular_font = hybrid.FontSpec(
            source_id="4",
            size=16.0,
            source_family="LMSans10-Regular",
            mapped_family="Arial",
            rgb=(0, 0, 0),
            family_bold=False,
            family_italic=False,
        )
        bold_font = hybrid.FontSpec(
            source_id="3",
            size=16.0,
            source_family="LMSans10-Bold",
            mapped_family="Arial",
            rgb=(0, 0, 0),
            family_bold=True,
            family_italic=False,
        )
        page, plan = plan_for(
            [
                span(7, "I", 55, 20, 9, font=BULLET_FONT),
                span(1, "Data/security:", 75, 20, 108, font=bold_font),
                span(2, "dependency and collective edges,", 190, 20, 220, font=regular_font),
                span(3, "signer/name/key-scope negative tests.", 75, 34.2, 263, font=regular_font),
                span(8, "I", 55, 48.4, 9, font=BULLET_FONT),
                span(4, "Environments:", 75, 48.4, 107, font=bold_font),
                span(5, "unit/integration harness first,", 190, 48.4, 210, font=regular_font),
                span(6, "container qualification.", 75, 62.6, 170, font=regular_font),
            ]
        )
        texts = [
            hybrid.expected_group_text(group)
            for group in plan.groups_by_page[page.number]
        ]
        self.assertIn(
            "Data/security: dependency and collective edges, "
            "signer/name/key-scope negative tests.",
            texts,
        )
        self.assertIn(
            "Environments: unit/integration harness first, "
            "container qualification.",
            texts,
        )

    def test_powerpoint_textbox_enables_automatic_word_wrap(self):
        page, plan = plan_for(
            [
                span(1, "A paragraph that wraps", 100, 20, 150),
                span(2, "inside one text box.", 100, 34.2, 120),
            ]
        )
        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        group = plan.groups_by_page[page.number][0]
        hybrid.add_text_group(slide, page, group)
        textbox = slide.shapes[-1]
        self.assertIs(textbox.text_frame.word_wrap, True)
        self.assertEqual(
            textbox.text, "A paragraph that wraps inside one text box."
        )

    def test_footer_fields_remain_independent_textboxes(self):
        page, plan = plan_for(
            [
                span(1, "Tianxing Ma", 11.4, 369.6, 67.6),
                span(2, "PhD Proposal Defense", 96.0, 369.6, 119.8),
                span(3, "University of Memphis", 232.7, 369.6, 118.7),
                span(4, "June 2026", 368.3, 369.6, 54.4),
                span(5, "32/62", 530.6, 369.6, 31.8),
            ]
        )
        texts = [
            hybrid.expected_group_text(group)
            for group in plan.groups_by_page[page.number]
        ]
        self.assertEqual(
            texts,
            [
                "Tianxing Ma",
                "PhD Proposal Defense",
                "University of Memphis",
                "June 2026",
                "32/62",
            ],
        )

    def test_frame_title_gets_full_slide_layout_width(self):
        title_font = hybrid.FontSpec(
            source_id="2",
            size=28.0,
            source_family="LMSans10-Regular",
            mapped_family="Arial",
            rgb=(26, 62, 122),
            family_bold=False,
            family_italic=False,
        )
        page, plan = plan_for(
            [
                span(
                    1,
                    "Four-Provider Mobility: Aggregate Results",
                    12.8,
                    14.8,
                    415.0,
                    font=title_font,
                )
            ]
        )
        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        hybrid.add_text_group(
            slide, page, plan.groups_by_page[page.number][0]
        )
        textbox = slide.shapes[-1]
        self.assertGreater(textbox.width / 914400, 15.0)
        self.assertIs(textbox.text_frame.word_wrap, False)
        source_left = hybrid.x_inches(12.8, page)
        textbox_left = textbox.left / 914400
        left_inset = textbox.text_frame.margin_left / 914400
        self.assertLess(textbox_left, source_left)
        self.assertGreater(left_inset, 0.0)
        self.assertAlmostEqual(textbox_left + left_inset, source_left, places=3)

    def test_cover_page_paragraphs_are_center_aligned(self):
        page, plan = plan_for(
            [
                span(1, "A centered title first line", 100, 20, 480),
                span(2, "short second line", 250, 34.2, 180),
                span(3, "Author Name", 290, 70, 100),
            ],
            page_number=1,
        )
        groups = plan.groups_by_page[page.number]
        self.assertTrue(groups)
        self.assertTrue(all(group.alignment == "center" for group in groups))

    def test_large_single_line_center_label_has_font_width_safety(self):
        large_font = hybrid.FontSpec(
            source_id="6",
            size=45.0,
            source_family="LMSans10-Regular",
            mapped_family="Arial",
            rgb=(0, 0, 0),
            family_bold=False,
            family_italic=False,
        )
        page, plan = plan_for(
            [span(1, "Questions?", 270, 220, 140, font=large_font)],
            page_number=51,
        )
        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        group = plan.groups_by_page[page.number][0]
        hybrid.add_text_group(slide, page, group)
        textbox = slide.shapes[-1]
        source_width = hybrid.x_inches(140, page)
        self.assertGreater(textbox.width / 914400, source_width + 0.4)
        self.assertEqual(textbox.text_frame.paragraphs[0].alignment, hybrid.PP_ALIGN.CENTER)

    def test_numbered_list_items_are_separate_wrapping_paragraphs(self):
        list_font = hybrid.FontSpec(
            source_id="5",
            size=16.0,
            source_family="LMSans10-Regular",
            mapped_family="Arial",
            rgb=(0, 0, 0),
            family_bold=False,
            family_italic=False,
        )
        page, plan = plan_for(
            [
                span(1, "1.", 54, 20, 13, font=list_font),
                span(2, "First roadmap item.", 75, 20, 145, font=list_font),
                span(3, "2.", 54, 34.2, 13, font=list_font),
                span(4, "Second roadmap item has a long first line", 75, 34.2, 300, font=list_font),
                span(5, "and a continuation line.", 75, 48.4, 170, font=list_font),
                span(6, "3.", 54, 62.6, 13, font=list_font),
                span(7, "Third roadmap item.", 75, 62.6, 150, font=list_font),
            ],
            page_number=4,
        )
        groups = plan.groups_by_page[page.number]
        texts = [hybrid.expected_group_text(group) for group in groups]
        self.assertEqual(
            texts,
            [
                "1. First roadmap item.",
                "2. Second roadmap item has a long first line and a continuation line.",
                "3. Third roadmap item.",
            ],
        )
        self.assertTrue(all(group.alignment == "left" for group in groups))

        presentation = Presentation()
        presentation.slide_width = Inches(hybrid.SLIDE_W_IN)
        presentation.slide_height = Inches(hybrid.SLIDE_H_IN)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        hybrid.add_text_group(slide, page, groups[1])
        paragraph = slide.shapes[-1].text_frame.paragraphs[0]
        paragraph_properties = paragraph._p.pPr
        self.assertGreater(int(paragraph_properties.get("marL")), 0)
        self.assertLess(int(paragraph_properties.get("indent")), 0)


if __name__ == "__main__":
    unittest.main()
